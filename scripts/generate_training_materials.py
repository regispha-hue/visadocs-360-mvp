#!/usr/bin/env python3
"""
VISADOCS Training Material Generator
Gera slides, vídeos e materiais de impressão para treinamento de POPs
Similar ao NotebookLM mas focado em materiais visuais de treinamento

Requisitos:
    pip install python-pptx moviepy fpdf pillow requests

Uso:
    python generate_training_materials.py --pop-id <id> --tenant-id <id>
    python generate_training_materials.py --pop-id POP001 --format all
"""

import argparse
import json
import os
import sys
from datetime import datetime
from typing import Dict, List, Optional
from dataclasses import dataclass
from pathlib import Path

# Core dependencies - serão instaladas automaticamente se não existirem
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print("Instalando python-pptx...")
    os.system("pip install python-pptx")
    from pptx import Presentation
    from pptx.util import Inches, Pt

try:
    from fpdf import FPDF
except ImportError:
    print("Instalando fpdf...")
    os.system("pip install fpdf")
    from fpdf import FPDF

try:
    from PIL import Image, ImageDraw, ImageFont
except ImportError:
    print("Instalando Pillow...")
    os.system("pip install Pillow")
    from PIL import Image, ImageDraw, ImageFont


@dataclass
class PopData:
    """Dados do POP para geração de materiais"""
    codigo: str
    titulo: str
    objetivo: str
    descricao: str
    setor: str
    responsavel: str
    versao: str
    equipe_envolvida: str = ""
    glossario: str = ""
    imagens_url: List[str] = None

    def __post_init__(self):
        if self.imagens_url is None:
            self.imagens_url = []


class MaterialGenerator:
    """Gerador de materiais visuais de treinamento"""
    
    # Cores VISADOCS
    COLORS = {
        'primary': RGBColor(13, 148, 136),      # Teal-600
        'secondary': RGBColor(15, 118, 110),    # Teal-700
        'accent': RGBColor(20, 184, 166),       # Teal-500
        'dark': RGBColor(30, 41, 59),           # Slate-800
        'light': RGBColor(248, 250, 252),       # Slate-50
        'white': RGBColor(255, 255, 255),
        'warning': RGBColor(234, 179, 8),      # Yellow-500
        'danger': RGBColor(239, 68, 68),        # Red-500
    }
    
    def __init__(self, output_dir: str = "training_materials"):
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(exist_ok=True)
        
    def generate_slide_deck(self, pop: PopData) -> str:
        """
        Gera apresentação PowerPoint do POP
        Retorna caminho do arquivo gerado
        """
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        # Slide 1: Capa
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)
        self._add_title_slide(slide, pop)
        
        # Slide 2: Objetivo
        slide = prs.slides.add_slide(slide_layout)
        self._add_content_slide(slide, "🎯 Objetivo", pop.objetivo)
        
        # Slide 3: Descrição
        slide = prs.slides.add_slide(slide_layout)
        self._add_content_slide(slide, "📝 Descrição do Procedimento", pop.descricao[:500] + "...")
        
        # Slide 4: Equipe
        if pop.equipe_envolvida:
            slide = prs.slides.add_slide(slide_layout)
            self._add_content_slide(slide, "👥 Equipe Envolvida", pop.equipe_envolvida)
        
        # Slide 5: Glossário
        if pop.glossario:
            slide = prs.slides.add_slide(slide_layout)
            self._add_content_slide(slide, "📚 Glossário", pop.glossario)
        
        # Slide 6: Responsáveis
        slide = prs.slides.add_slide(slide_layout)
        self._add_responsibles_slide(slide, pop)
        
        # Slide 7: QR Code para mais info
        slide = prs.slides.add_slide(slide_layout)
        self._add_qr_slide(slide, pop)
        
        # Salvar
        filename = f"{pop.codigo}_Slides_{datetime.now().strftime('%Y%m%d')}.pptx"
        filepath = self.output_dir / filename
        prs.save(filepath)
        print(f"✅ Slides gerados: {filepath}")
        return str(filepath)
    
    def _add_title_slide(self, slide, pop: PopData):
        """Adiciona slide de título"""
        # Background shape
        shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.COLORS['primary']
        shape.line.fill.background()
        
        # Título
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = pop.titulo
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = self.COLORS['white']
        p.alignment = PP_ALIGN.CENTER
        
        # Código
        code_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12), Inches(0.8))
        tf = code_box.text_frame
        p = tf.paragraphs[0]
        p.text = pop.codigo
        p.font.size = Pt(28)
        p.font.color.rgb = self.COLORS['light']
        p.alignment = PP_ALIGN.CENTER
        
        # Setor badge
        badge_box = slide.shapes.add_textbox(Inches(5), Inches(4.5), Inches(3.333), Inches(0.6))
        tf = badge_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"  📍 {pop.setor}  "
        p.font.size = Pt(18)
        p.font.color.rgb = self.COLORS['white']
        p.alignment = PP_ALIGN.CENTER
    
    def _add_content_slide(self, slide, title: str, content: str):
        """Adiciona slide de conteúdo"""
        # Header bar
        shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.COLORS['primary']
        shape.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.COLORS['white']
        
        # Content
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12), Inches(5))
        tf = content_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = content
        p.font.size = Pt(20)
        p.font.color.rgb = self.COLORS['dark']
        p.line_spacing = 1.5
    
    def _add_responsibles_slide(self, slide, pop: PopData):
        """Slide de responsáveis"""
        self._add_content_slide(slide, "✅ Responsáveis", 
            f"Responsável: {pop.responsavel}\n\n"
            f"Versão: {pop.versao}\n\n"
            f"Setor: {pop.setor}")
    
    def _add_qr_slide(self, slide, pop: PopData):
        """Slide com QR code"""
        self._add_content_slide(slide, "📱 Acesse o POP Completo", 
            "Escaneie o QR code ou acesse:\n"
            f"visadocs.com/pop/{pop.codigo}\n\n"
            "Tenha acesso ao procedimento completo,\n"
            "documentos relacionados e certificados.")
    
    def generate_poster(self, pop: PopData, size: str = "A3") -> str:
        """
        Gera poster para impressão (como os de lavagem das mãos)
        Tamanhos: A4, A3, A2
        """
        # Tamanhos em pixels (300 DPI)
        sizes = {
            'A4': (2480, 3508),   # 210 x 297 mm
            'A3': (3508, 4961),   # 297 x 420 mm
            'A2': (4961, 7016),   # 420 x 594 mm
        }
        
        width, height = sizes.get(size, sizes['A3'])
        
        # Criar imagem
        img = Image.new('RGB', (width, height), 'white')
        draw = ImageDraw.Draw(img)
        
        # Cores
        teal = (13, 148, 136)
        dark = (30, 41, 59)
        light = (248, 250, 252)
        
        try:
            font_title = ImageFont.truetype("arial.ttf", 80)
            font_heading = ImageFont.truetype("arial.ttf", 50)
            font_body = ImageFont.truetype("arial.ttf", 36)
        except:
            font_title = ImageFont.load_default()
            font_heading = font_title
            font_body = font_title
        
        # Header
        draw.rectangle([0, 0, width, 300], fill=teal)
        draw.text((60, 100), pop.codigo, fill='white', font=font_heading)
        
        # Título
        y_pos = 350
        draw.text((60, y_pos), pop.titulo[:80], fill=dark, font=font_title)
        y_pos += 150
        
        # Objetivo
        draw.text((60, y_pos), "🎯 OBJETIVO", fill=teal, font=font_heading)
        y_pos += 80
        
        # Quebrar texto em linhas
        words = pop.objetivo.split()
        lines = []
        current_line = []
        for word in words:
            current_line.append(word)
            if len(' '.join(current_line)) > 60:
                lines.append(' '.join(current_line[:-1]))
                current_line = [current_line[-1]]
        if current_line:
            lines.append(' '.join(current_line))
        
        for line in lines[:8]:  # Max 8 linhas
            draw.text((80, y_pos), line, fill=dark, font=font_body)
            y_pos += 50
        
        # Procedimento em passos (simplificado)
        y_pos += 50
        draw.text((60, y_pos), "📝 PROCEDIMENTO", fill=teal, font=font_heading)
        y_pos += 80
        
        # Extrair passos da descrição
        passos = pop.descricao.split('.')[:5]
        for i, passo in enumerate(passos, 1):
            if len(passo.strip()) > 10:
                draw.text((80, y_pos), f"{i}. {passo.strip()[:100]}...", fill=dark, font=font_body)
                y_pos += 60
        
        # Responsável
        y_pos = height - 200
        draw.rectangle([0, y_pos-20, width, height], fill=light)
        draw.text((60, y_pos), f"✅ Responsável: {pop.responsavel}", fill=dark, font=font_body)
        draw.text((60, y_pos+60), f"📍 {pop.setor} | Versão: {pop.versao}", fill='gray', font=font_body)
        
        # Salvar
        filename = f"{pop.codigo}_Poster_{size}_{datetime.now().strftime('%Y%m%d')}.png"
        filepath = self.output_dir / filename
        img.save(filepath, 'PNG', dpi=(300, 300))
        print(f"✅ Poster gerado: {filepath}")
        return str(filepath)
    
    def generate_pdf_guide(self, pop: PopData) -> str:
        """
        Gera guia em PDF para impressão (formato caderno)
        """
        pdf = FPDF()
        pdf.add_page()
        
        # Configurar fontes
        pdf.set_auto_page_break(auto=True, margin=15)
        
        # Header
        pdf.set_fill_color(13, 148, 136)
        pdf.rect(0, 0, 210, 30, 'F')
        pdf.set_font('Arial', 'B', 12)
        pdf.set_text_color(255, 255, 255)
        pdf.cell(0, 10, f"{pop.codigo} - {pop.titulo[:50]}", ln=True, align='C')
        
        # Conteúdo
        pdf.set_y(40)
        pdf.set_text_color(30, 41, 59)
        pdf.set_font('Arial', 'B', 16)
        pdf.cell(0, 10, pop.titulo, ln=True)
        
        pdf.set_font('Arial', '', 11)
        pdf.ln(5)
        
        # Objetivo
        pdf.set_font('Arial', 'B', 12)
        pdf.cell(0, 8, "OBJETIVO", ln=True)
        pdf.set_font('Arial', '', 11)
        pdf.multi_cell(0, 6, pop.objetivo)
        pdf.ln(5)
        
        # Descrição
        pdf.set_font('Arial', 'B', 12)
        pdf.cell(0, 8, "DESCRI\xc7\xc3O DO PROCEDIMENTO", ln=True)
        pdf.set_font('Arial', '', 11)
        
        # Dividir em parágrafos
        descricao = pop.descricao[:2000]  # Limitar tamanho
        paragraphs = descricao.split('\n')
        for para in paragraphs:
            if para.strip():
                pdf.multi_cell(0, 6, para.strip())
                pdf.ln(2)
        
        # Responsáveis
        pdf.ln(5)
        pdf.set_font('Arial', 'B', 12)
        pdf.cell(0, 8, "RESPONS\xc1VEIS", ln=True)
        pdf.set_font('Arial', '', 11)
        pdf.cell(0, 6, f"Respons\xe1vel: {pop.responsavel}", ln=True)
        pdf.cell(0, 6, f"Setor: {pop.setor}", ln=True)
        pdf.cell(0, 6, f"Vers\xe3o: {pop.versao}", ln=True)
        
        # Salvar
        filename = f"{pop.codigo}_Guia_{datetime.now().strftime('%Y%m%d')}.pdf"
        filepath = self.output_dir / filename
        pdf.output(str(filepath))
        print(f"✅ Guia PDF gerado: {filepath}")
        return str(filepath)
    
    def generate_video_script(self, pop: PopData) -> str:
        """
        Gera roteiro para vídeo de treinamento
        """
        script = f"""
ROTEIRO DE VÍDEO DE TREINAMENTO
{pop.codigo} - {pop.titulo}
Gerado em: {datetime.now().strftime('%d/%m/%Y %H:%M')}
{'='*60}

DURAÇÃO ESTIMADA: 3-5 minutos
FORMATO: Vídeo institucional + Motion Graphics

CENA 1 - INTRODUÇÃO (0:00-0:30)
--------------------------------
[Visual]: Logo VISADOCS + Código do POP
[Texto em tela]: "{pop.codigo}"
[Áudio]: "Bem-vindo ao treinamento do procedimento {pop.titulo}"

CENA 2 - CONTEXTO (0:30-1:00)
------------------------------
[Visual]: Setor {pop.setor} em operação
[Texto em tela]: "Setor: {pop.setor}"
[Áudio]: "Este procedimento é aplicável no setor de {pop.setor}"

CENA 3 - OBJETIVO (1:00-1:30)
-----------------------------
[Visual]: Ícones animados ilustrando o objetivo
[Texto em tela]: "Objetivo"
[Áudio]: "{pop.objetivo[:150]}..."

CENA 4 - PROCEDIMENTO (1:30-3:30)
----------------------------------
[Visual]: Animação passo a passo
"""
        # Adicionar passos
        passos = pop.descricao.split('.')[:8]
        for i, passo in enumerate(passos, 1):
            if len(passo.strip()) > 10:
                script += f"""
Passo {i}:
[Visual]: Ilustração do passo {i}
[Texto em tela]: "{passo.strip()[:100]}"
[Áudio]: "Passo {i}: {passo.strip()[:100]}"
"""
        
        script += f"""
CENA 5 - RESPONSÁVEIS (3:30-4:00)
----------------------------------
[Visual]: Foto ou avatar do responsável
[Texto em tela]: "Responsável: {pop.responsavel}"
[Áudio]: "Em caso de dúvidas, procure {pop.responsavel}"

CENA 6 - ENCERRAMENTO (4:00-5:00)
----------------------------------
[Visual]: QR Code para acesso ao POP completo
[Texto em tela]: "Escaneie para mais informações"
[Áudio]: "Para mais detalhes, acesse o POP completo no sistema VISADOCS"

MATERIAIS NECESSÁRIOS:
- Ilustrações dos procedimentos
- Fotos do setor
- Motion graphics
- Narração profissional
- Trilha sonora institucional

DICAS DE PRODUÇÃO:
- Manter linguagem simples e clara
- Usar legendas para acessibilidade
- Incluir pausas entre passos
- Destacar pontos críticos em vermelho
"""
        
        filename = f"{pop.codigo}_Roteiro_Video_{datetime.now().strftime('%Y%m%d')}.txt"
        filepath = self.output_dir / filename
        with open(filepath, 'w', encoding='utf-8') as f:
            f.write(script)
        print(f"✅ Roteiro de vídeo gerado: {filepath}")
        return str(filepath)
    
    def generate_all(self, pop: PopData) -> List[str]:
        """Gera todos os materiais de uma vez"""
        print(f"\n🎬 Gerando materiais de treinamento para: {pop.codigo}")
        print("="*60)
        
        files = []
        
        # Slides
        try:
            files.append(self.generate_slide_deck(pop))
        except Exception as e:
            print(f"⚠️ Erro ao gerar slides: {e}")
        
        # Poster A3
        try:
            files.append(self.generate_poster(pop, "A3"))
        except Exception as e:
            print(f"⚠️ Erro ao gerar poster: {e}")
        
        # PDF
        try:
            files.append(self.generate_pdf_guide(pop))
        except Exception as e:
            print(f"⚠️ Erro ao gerar PDF: {e}")
        
        # Roteiro
        try:
            files.append(self.generate_video_script(pop))
        except Exception as e:
            print(f"⚠️ Erro ao gerar roteiro: {e}")
        
        print("\n" + "="*60)
        print(f"✅ Total de arquivos gerados: {len(files)}")
        print(f"📁 Diretório: {self.output_dir.absolute()}")
        
        return files


def mock_pop_data() -> PopData:
    """Dados mockados para teste"""
    return PopData(
        codigo="POP.001",
        titulo="Lavagem das Mãos em Farmácia",
        objetivo="Estabelecer o procedimento correto para lavagem das mãos visando garantir a higiene e prevenir contaminação cruzada na manipulação de medicamentos.",
        descricao="""
1. Molhar as mãos e punhos com água corrente.
2. Aplicar sabonete antisséptico suficiente para cobrir todas as superfícies.
3. Esfregar palmas entre si com movimentos circulares.
4. Esfregar dorso de uma mão contra palma da outra, alternando.
5. Entrelacear dedos e esfregar entre eles.
6. Esfregar polegares envolvendo-os com a outra mão.
7. Esfregar unhas contra as palmas.
8. Enxaguar abundantemente com água corrente.
9. Secar com papel toalha descartável.
10. Desligar torneira usando o cotovelo ou papel.
        """.strip(),
        setor="Manipulação",
        responsavel="Farmacêutico Responsável",
        versao="Rev.03",
        equipe_envolvida="Todos os colaboradores do setor de manipulação",
        glossario="Antisséptico: Substância que inibe o crescimento de microrganismos na pele."
    )


def main():
    parser = argparse.ArgumentParser(
        description="VISADOCS Training Material Generator",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Exemplos:
  python generate_training_materials.py --mock
  python generate_training_materials.py --pop-id POP001 --format all
  python generate_training_materials.py --pop-id POP001 --format slides
        """
    )
    
    parser.add_argument("--mock", action="store_true", 
                       help="Usar dados mockados para teste")
    parser.add_argument("--pop-id", type=str,
                       help="ID do POP no sistema")
    parser.add_argument("--format", type=str, 
                       choices=["all", "slides", "poster", "pdf", "script"],
                       default="all",
                       help="Formato de saída")
    parser.add_argument("--output", type=str, 
                       default="training_materials",
                       help="Diretório de saída")
    
    args = parser.parse_args()
    
    # Criar gerador
    generator = MaterialGenerator(args.output)
    
    # Obter dados do POP
    if args.mock:
        pop = mock_pop_data()
    else:
        # Aqui integraria com API do VISADOCS
        print("⚠️ Modo mock ativado (use --mock explicitamente)")
        pop = mock_pop_data()
    
    # Gerar materiais
    if args.format == "all":
        generator.generate_all(pop)
    elif args.format == "slides":
        generator.generate_slide_deck(pop)
    elif args.format == "poster":
        generator.generate_poster(pop, "A3")
    elif args.format == "pdf":
        generator.generate_pdf_guide(pop)
    elif args.format == "script":
        generator.generate_video_script(pop)
    
    print("\n🎉 Materiais gerados com sucesso!")


if __name__ == "__main__":
    main()
